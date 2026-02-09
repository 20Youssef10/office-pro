"""
Presentation Module
Professional presentation editing with support for PPTX, PPT, and ODP files.
"""

import os
from pathlib import Path
from typing import List, Optional

try:
    from PyQt6.QtWidgets import (
        QWidget,
        QVBoxLayout,
        QHBoxLayout,
        QGraphicsView,
        QGraphicsScene,
        QGraphicsRectItem,
        QGraphicsTextItem,
        QToolBar,
        QPushButton,
        QComboBox,
        QFileDialog,
        QMessageBox,
        QInputDialog,
        QListWidget,
        QListWidgetItem,
        QLabel,
        QFrame,
        QSplitter,
        QTextEdit,
        QLineEdit,
        QColorDialog,
        QFontDialog,
        QGraphicsItem,
    )
    from PyQt6.QtCore import Qt, QSize, QRectF, pyqtSignal
    from PyQt6.QtGui import QPen, QBrush, QColor, QFont, QAction, QPainter, QPixmap

    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.enum.text import PP_ALIGN
        from pptx.dml.color import RGBColor

        PPTX_AVAILABLE = True
    except ImportError:
        PPTX_AVAILABLE = False

except ImportError as e:
    print(f"Import error in presentation: {e}")
    raise


class SlideCanvas(QGraphicsView):
    """Canvas for editing slides"""

    item_selected = pyqtSignal(object)

    def __init__(self, parent=None):
        super().__init__(parent)
        self.scene = QGraphicsScene(self)
        self.setScene(self.scene)
        self.setup_canvas()

        # Slide dimensions (16:9 aspect ratio)
        self.slide_width = 960
        self.slide_height = 540

        # Current slide background
        self.background = None

        self.create_slide_background()

    def setup_canvas(self):
        """Setup canvas properties"""
        self.setStyleSheet("""
            QGraphicsView {
                background-color: #e0e0e0;
                border: none;
            }
        """)
        self.setRenderHint(QPainter.RenderHint.Antialiasing)
        self.setViewportUpdateMode(QGraphicsView.ViewportUpdateMode.FullViewportUpdate)
        self.setHorizontalScrollBarPolicy(Qt.ScrollBarPolicy.ScrollBarAsNeeded)
        self.setVerticalScrollBarPolicy(Qt.ScrollBarPolicy.ScrollBarAsNeeded)
        self.setDragMode(QGraphicsView.DragMode.RubberBandDrag)

    def create_slide_background(self):
        """Create slide background"""
        if self.background:
            self.scene.removeItem(self.background)

        self.background = QGraphicsRectItem(0, 0, self.slide_width, self.slide_height)
        self.background.setBrush(QBrush(QColor("white")))
        self.background.setPen(QPen(QColor("#cccccc"), 1))
        self.scene.addItem(self.background)

        self.setSceneRect(-50, -50, self.slide_width + 100, self.slide_height + 100)

    def add_text_box(self, x=50, y=50, width=400, height=100):
        """Add a text box to the slide"""
        text_item = QGraphicsTextItem()
        text_item.setPos(x, y)
        text_item.setTextWidth(width)
        text_item.setPlainText("Click to edit")
        text_item.setDefaultTextColor(QColor("black"))
        text_item.setFont(QFont("Calibri", 24))
        text_item.setFlags(
            QGraphicsItem.GraphicsItemFlag.ItemIsMovable
            | QGraphicsItem.GraphicsItemFlag.ItemIsSelectable
        )
        self.scene.addItem(text_item)
        return text_item

    def add_shape(self, shape_type="rectangle", x=50, y=50, width=200, height=150):
        """Add a shape to the slide"""
        if shape_type == "rectangle":
            shape = QGraphicsRectItem(x, y, width, height)
            shape.setBrush(QBrush(QColor("#4472C4")))
            shape.setPen(QPen(QColor("#2F5597"), 2))
        else:
            shape = QGraphicsRectItem(x, y, width, height)
            shape.setBrush(QBrush(QColor("#70AD47")))
            shape.setPen(QPen(QColor("#507E32"), 2))

        shape.setFlags(
            QGraphicsItem.GraphicsItemFlag.ItemIsMovable
            | QGraphicsItem.GraphicsItemFlag.ItemIsSelectable
        )
        self.scene.addItem(shape)
        return shape

    def clear_slide(self):
        """Clear all items except background"""
        items_to_remove = [
            item for item in self.scene.items() if item != self.background
        ]
        for item in items_to_remove:
            self.scene.removeItem(item)

    def get_slide_data(self):
        """Get slide data for saving"""
        items_data = []
        for item in self.scene.items():
            if item != self.background:
                item_data = {
                    "type": type(item).__name__,
                    "pos": (item.pos().x(), item.pos().y()),
                }

                if isinstance(item, QGraphicsTextItem):
                    item_data["text"] = item.toPlainText()
                    item_data["font"] = item.font().toString()
                    item_data["color"] = item.defaultTextColor().name()
                    item_data["width"] = item.textWidth()
                elif isinstance(item, QGraphicsRectItem):
                    rect = item.rect()
                    item_data["rect"] = (
                        rect.x(),
                        rect.y(),
                        rect.width(),
                        rect.height(),
                    )
                    item_data["brush"] = item.brush().color().name()
                    item_data["pen"] = item.pen().color().name()

                items_data.append(item_data)

        return items_data

    def load_slide_data(self, items_data):
        """Load slide data"""
        self.clear_slide()

        for item_data in items_data:
            item_type = item_data.get("type")

            if item_type == "QGraphicsTextItem":
                text_item = self.add_text_box(
                    item_data["pos"][0],
                    item_data["pos"][1],
                    item_data.get("width", 400),
                    100,
                )
                text_item.setPlainText(item_data.get("text", ""))

                font = QFont()
                font.fromString(item_data.get("font", "Calibri,24,-1,5,50,0,0,0,0,0"))
                text_item.setFont(font)
                text_item.setDefaultTextColor(QColor(item_data.get("color", "#000000")))

            elif item_type == "QGraphicsRectItem":
                rect = item_data.get("rect", (50, 50, 200, 150))
                shape = self.add_shape("rectangle", rect[0], rect[1], rect[2], rect[3])
                shape.setBrush(QBrush(QColor(item_data.get("brush", "#4472C4"))))
                shape.setPen(QPen(QColor(item_data.get("pen", "#2F5597")), 2))


class PresentationEditor(QWidget):
    """Professional Presentation Editor"""

    def __init__(self, parent=None):
        super().__init__(parent)
        self.parent = parent
        self.current_file = None
        self.is_modified = False
        self.slides = []
        self.current_slide_index = 0

        self.init_ui()
        self.setup_toolbar()
        self.new_presentation()

    def init_ui(self):
        """Initialize the user interface"""
        layout = QVBoxLayout(self)
        layout.setContentsMargins(0, 0, 0, 0)
        layout.setSpacing(0)

        # Toolbar
        self.toolbar = QToolBar()
        self.toolbar.setMovable(False)
        self.toolbar.setStyleSheet("""
            QToolBar {
                background-color: #d24726;
                border: none;
                padding: 5px;
                spacing: 5px;
            }
            QToolBar QPushButton {
                background-color: rgba(255,255,255,0.1);
                color: white;
                border: none;
                padding: 8px 15px;
                border-radius: 4px;
                font-size: 12px;
            }
            QToolBar QPushButton:hover {
                background-color: rgba(255,255,255,0.2);
            }
        """)
        layout.addWidget(self.toolbar)

        # Main content area
        splitter = QSplitter(Qt.Orientation.Horizontal)

        # Slide panel (left)
        slide_panel = QFrame()
        slide_panel.setStyleSheet("""
            QFrame {
                background-color: #f5f5f5;
                border-right: 1px solid #ddd;
                min-width: 200px;
                max-width: 250px;
            }
        """)
        slide_layout = QVBoxLayout(slide_panel)
        slide_layout.setContentsMargins(5, 5, 5, 5)

        slide_label = QLabel("Slides")
        slide_label.setStyleSheet("font-weight: bold; font-size: 14px;")
        slide_layout.addWidget(slide_label)

        self.slide_list = QListWidget()
        self.slide_list.setStyleSheet("""
            QListWidget {
                background-color: white;
                border: 1px solid #ddd;
            }
            QListWidget::item {
                padding: 10px;
                border-bottom: 1px solid #eee;
            }
            QListWidget::item:selected {
                background-color: #d24726;
                color: white;
            }
        """)
        self.slide_list.currentRowChanged.connect(self.on_slide_selected)
        slide_layout.addWidget(self.slide_list)

        # Slide panel buttons
        btn_layout = QHBoxLayout()

        self.add_slide_btn = QPushButton("+ Slide")
        self.add_slide_btn.clicked.connect(self.add_slide)
        btn_layout.addWidget(self.add_slide_btn)

        self.del_slide_btn = QPushButton("- Slide")
        self.del_slide_btn.clicked.connect(self.delete_slide)
        btn_layout.addWidget(self.del_slide_btn)

        slide_layout.addLayout(btn_layout)

        splitter.addWidget(slide_panel)

        # Canvas (center)
        self.canvas = SlideCanvas(self)
        splitter.addWidget(self.canvas)

        splitter.setSizes([220, 1180])
        layout.addWidget(splitter)

        # Status bar
        self.status_frame = QFrame()
        self.status_frame.setStyleSheet("""
            QFrame {
                background-color: #f0f0f0;
                border-top: 1px solid #ddd;
                padding: 5px;
            }
        """)
        status_layout = QHBoxLayout(self.status_frame)
        status_layout.setContentsMargins(10, 5, 10, 5)

        self.slide_count_label = QLabel("Slides: 1")
        self.current_slide_label = QLabel("Current: 1/1")

        status_layout.addWidget(self.slide_count_label)
        status_layout.addStretch()
        status_layout.addWidget(self.current_slide_label)

        layout.addWidget(self.status_frame)

    def setup_toolbar(self):
        """Setup the toolbar"""
        # File operations
        self.new_btn = QPushButton("New")
        self.new_btn.clicked.connect(self.new_presentation)
        self.toolbar.addWidget(self.new_btn)

        self.open_btn = QPushButton("Open")
        self.open_btn.clicked.connect(self.open_file_dialog)
        self.toolbar.addWidget(self.open_btn)

        self.save_btn = QPushButton("Save")
        self.save_btn.clicked.connect(self.save)
        self.toolbar.addWidget(self.save_btn)

        self.toolbar.addSeparator()

        # Insert
        self.add_text_btn = QPushButton("Text Box")
        self.add_text_btn.clicked.connect(self.add_text_box)
        self.toolbar.addWidget(self.add_text_btn)

        self.add_shape_btn = QPushButton("Shape")
        self.add_shape_btn.clicked.connect(self.add_shape)
        self.toolbar.addWidget(self.add_shape_btn)

        self.toolbar.addSeparator()

        # Navigation
        self.prev_btn = QPushButton("← Prev")
        self.prev_btn.clicked.connect(self.previous_slide)
        self.toolbar.addWidget(self.prev_btn)

        self.next_btn = QPushButton("Next →")
        self.next_btn.clicked.connect(self.next_slide)
        self.toolbar.addWidget(self.next_btn)

        self.toolbar.addSeparator()

        # Home button
        self.home_btn = QPushButton("← Home")
        self.home_btn.clicked.connect(self.go_home)
        self.toolbar.addWidget(self.home_btn)

    def new_presentation(self):
        """Create new presentation"""
        if self.check_save():
            self.slides = [{"items": []}]
            self.current_slide_index = 0
            self.current_file = None
            self.is_modified = False

            self.canvas.clear_slide()
            self.update_slide_list()
            self.update_title()
            self.update_status()

    def add_slide(self):
        """Add new slide"""
        # Save current slide data
        self.save_current_slide()

        # Add new slide
        self.slides.append({"items": []})
        self.current_slide_index = len(self.slides) - 1

        # Update UI
        self.update_slide_list()
        self.slide_list.setCurrentRow(self.current_slide_index)
        self.canvas.clear_slide()
        self.is_modified = True
        self.update_status()

    def delete_slide(self):
        """Delete current slide"""
        if len(self.slides) > 1:
            self.slides.pop(self.current_slide_index)

            if self.current_slide_index >= len(self.slides):
                self.current_slide_index = len(self.slides) - 1

            self.update_slide_list()
            self.slide_list.setCurrentRow(self.current_slide_index)
            self.load_current_slide()
            self.is_modified = True
            self.update_status()
        else:
            QMessageBox.warning(
                self, "Cannot Delete", "You must have at least one slide."
            )

    def on_slide_selected(self, index):
        """Handle slide selection change"""
        if index >= 0 and index < len(self.slides):
            self.save_current_slide()
            self.current_slide_index = index
            self.load_current_slide()
            self.update_status()

    def save_current_slide(self):
        """Save current slide data"""
        if 0 <= self.current_slide_index < len(self.slides):
            self.slides[self.current_slide_index]["items"] = (
                self.canvas.get_slide_data()
            )

    def load_current_slide(self):
        """Load current slide data"""
        if 0 <= self.current_slide_index < len(self.slides):
            items_data = self.slides[self.current_slide_index].get("items", [])
            self.canvas.load_slide_data(items_data)

    def update_slide_list(self):
        """Update slide list widget"""
        self.slide_list.clear()
        for i, slide in enumerate(self.slides):
            item = QListWidgetItem(f"Slide {i + 1}")
            self.slide_list.addItem(item)

    def previous_slide(self):
        """Go to previous slide"""
        if self.current_slide_index > 0:
            self.slide_list.setCurrentRow(self.current_slide_index - 1)

    def next_slide(self):
        """Go to next slide"""
        if self.current_slide_index < len(self.slides) - 1:
            self.slide_list.setCurrentRow(self.current_slide_index + 1)

    def add_text_box(self):
        """Add text box to current slide"""
        self.canvas.add_text_box()
        self.is_modified = True
        self.update_title()

    def add_shape(self):
        """Add shape to current slide"""
        shapes = ["Rectangle", "Ellipse", "Rounded Rectangle"]
        shape, ok = QInputDialog.getItem(
            self, "Add Shape", "Shape type:", shapes, 0, False
        )
        if ok:
            shape_type = shape.lower().replace(" ", "_")
            self.canvas.add_shape(shape_type)
            self.is_modified = True
            self.update_title()

    def open_file_dialog(self):
        """Open file dialog"""
        if self.check_save():
            file_path, _ = QFileDialog.getOpenFileName(
                self,
                "Open Presentation",
                "",
                "PowerPoint (*.pptx *.ppt);;All Files (*)",
            )
            if file_path:
                self.open_presentation(file_path)

    def open_presentation(self, file_path):
        """Open presentation file"""
        try:
            if file_path.endswith(".pptx") and PPTX_AVAILABLE:
                self.open_pptx(file_path)
            else:
                QMessageBox.warning(
                    self,
                    "Error",
                    "PPTX support requires python-pptx. Install it first.",
                )
                return

            self.current_file = file_path
            self.is_modified = False
            self.update_title()
            self.update_slide_list()
            self.slide_list.setCurrentRow(0)
            self.load_current_slide()
            self.update_status()

        except Exception as e:
            QMessageBox.critical(
                self, "Error", f"Failed to open presentation:\n{str(e)}"
            )

    def open_pptx(self, file_path):
        """Open PowerPoint file"""
        prs = Presentation(file_path)

        self.slides = []
        for slide in prs.slides:
            slide_data = {"items": []}

            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    item_data = {
                        "type": "QGraphicsTextItem",
                        "text": shape.text,
                        "pos": (shape.left / 914400 * 96, shape.top / 914400 * 96),
                        "width": 400,
                        "font": "Calibri,24,-1,5,50,0,0,0,0,0",
                        "color": "#000000",
                    }
                    slide_data["items"].append(item_data)

            self.slides.append(slide_data)

        if not self.slides:
            self.slides = [{"items": []}]

    def save(self):
        """Save presentation"""
        # Save current slide first
        self.save_current_slide()

        if self.current_file:
            self.save_as(self.current_file)
        else:
            self.save_as_dialog()

    def save_as_dialog(self):
        """Save as dialog"""
        file_path, _ = QFileDialog.getSaveFileName(
            self, "Save Presentation", "", "PowerPoint (*.pptx);;All Files (*)"
        )
        if file_path:
            self.save_as(file_path)

    def save_as(self, file_path):
        """Save presentation to file"""
        try:
            if file_path.endswith(".pptx") and PPTX_AVAILABLE:
                self.save_as_pptx(file_path)
            else:
                QMessageBox.warning(
                    self, "Error", "Can only save as PPTX. Install python-pptx."
                )
                return

            self.current_file = file_path
            self.is_modified = False
            self.update_title()

            if self.parent:
                self.parent.status_bar.showMessage(f"Saved: {file_path}", 3000)

        except Exception as e:
            QMessageBox.critical(
                self, "Error", f"Failed to save presentation:\n{str(e)}"
            )

    def save_as_pptx(self, file_path):
        """Save as PowerPoint file"""
        prs = Presentation()

        for slide_data in self.slides:
            slide_layout = prs.slide_layouts[6]  # Blank layout
            slide = prs.slides.add_slide(slide_layout)

            for item in slide_data.get("items", []):
                if item.get("type") == "QGraphicsTextItem":
                    left = Inches(item["pos"][0] / 96)
                    top = Inches(item["pos"][1] / 96)
                    width = Inches(4)
                    height = Inches(1)

                    textbox = slide.shapes.add_textbox(left, top, width, height)
                    text_frame = textbox.text_frame
                    text_frame.text = item.get("text", "")

        prs.save(file_path)

    def check_save(self):
        """Check if save is needed"""
        if self.is_modified:
            reply = QMessageBox.question(
                self,
                "Unsaved Changes",
                "The presentation has unsaved changes. Save now?",
                QMessageBox.StandardButton.Save
                | QMessageBox.StandardButton.Discard
                | QMessageBox.StandardButton.Cancel,
            )

            if reply == QMessageBox.StandardButton.Save:
                self.save()
                return True
            elif reply == QMessageBox.StandardButton.Cancel:
                return False

        return True

    def update_title(self):
        """Update window title"""
        if self.parent:
            filename = "Untitled"
            if self.current_file:
                filename = Path(self.current_file).name
            modified = "*" if self.is_modified else ""
            self.parent.setWindowTitle(
                f"Office Pro - Presentation - {filename}{modified}"
            )

    def update_status(self):
        """Update status bar"""
        self.slide_count_label.setText(f"Slides: {len(self.slides)}")
        self.current_slide_label.setText(
            f"Current: {self.current_slide_index + 1}/{len(self.slides)}"
        )

    def go_home(self):
        """Return to main menu"""
        self.save_current_slide()
        if self.check_save():
            if self.parent:
                self.parent.show_main_menu()
